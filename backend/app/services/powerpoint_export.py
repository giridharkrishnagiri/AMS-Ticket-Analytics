from __future__ import annotations

import html
import re
from datetime import UTC, datetime
from io import BytesIO
from typing import Any
from uuid import UUID

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.presentation import Presentation as PresentationType
from pptx.slide import Slide
from pptx.util import Inches, Pt
from sqlalchemy.orm import Session

from app.services.offline_dashboard_export import build_offline_dashboard_payload

FunctionalFilter = str | frozenset[str]

TEAL = RGBColor(15, 118, 110)
BLUE = RGBColor(37, 99, 235)
ORANGE = RGBColor(234, 88, 12)
RED = RGBColor(185, 28, 28)
PURPLE = RGBColor(124, 58, 237)
GREEN = RGBColor(22, 163, 74)
SLATE = RGBColor(71, 85, 105)
TEXT = RGBColor(17, 24, 39)
MUTED = RGBColor(82, 98, 122)
BORDER = RGBColor(217, 226, 236)
LIGHT_BG = RGBColor(248, 250, 252)
WHITE = RGBColor(255, 255, 255)

CHART_COLORS = [TEAL, BLUE, ORANGE, PURPLE, GREEN, RED, SLATE]
PRIORITY_COLORS = [BLUE, ORANGE, TEAL, PURPLE]
DURATION_BUCKETS = ("0-1 day", "1-3 days", "3-10 days", ">10 days")
MTTR_PRIORITIES = ("P1", "P2", "P3", "P4")


def powerpoint_filename(exported_at: datetime) -> str:
    return f"AMS_Apps_Volumetrics_Dashboard_{exported_at.strftime('%Y%m%d_%H%M')}.pptx"


def safe_text(value: Any) -> str:
    return "" if value is None else str(value)


def strip_html(value: str | None) -> str:
    if not value:
        return ""
    without_tags = re.sub(r"<br\s*/?>", "\n", value, flags=re.IGNORECASE)
    without_tags = re.sub(r"</p\s*>", "\n", without_tags, flags=re.IGNORECASE)
    without_tags = re.sub(r"<[^>]+>", "", without_tags)
    return html.unescape(without_tags).strip()


def fmt_number(value: Any, digits: int = 0) -> str:
    try:
        number = float(value or 0)
    except (TypeError, ValueError):
        return "0"
    if digits:
        return f"{number:,.{digits}f}"
    return f"{number:,.0f}"


def fmt_pct(value: Any) -> str:
    if value is None:
        return "N/A"
    try:
        return f"{float(value):.1f}%"
    except (TypeError, ValueError):
        return "N/A"


def add_blank_slide(presentation: PresentationType) -> Slide:
    return presentation.slides.add_slide(presentation.slide_layouts[6])


def set_font(run: Any, *, size: int, bold: bool = False, color: RGBColor = TEXT) -> None:
    run.font.name = "Aptos"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def add_textbox(
    slide: Slide,
    text: str,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    size: int = 14,
    bold: bool = False,
    color: RGBColor = TEXT,
    align: PP_ALIGN = PP_ALIGN.LEFT,
) -> Any:
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.vertical_anchor = MSO_ANCHOR.TOP
    for index, line in enumerate(text.splitlines() or [""]):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.alignment = align
        run = paragraph.add_run()
        run.text = line
        set_font(run, size=size, bold=bold, color=color)
    return shape


def add_header(slide: Slide, title: str, subtitle: str | None = None) -> None:
    add_textbox(slide, title, 0.45, 0.25, 12.2, 0.45, size=22, bold=True, color=TEAL)
    if subtitle:
        add_textbox(slide, subtitle, 0.47, 0.72, 12.1, 0.28, size=9, color=MUTED)
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.45), Inches(1.02), Inches(12.35), Inches(0.02)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = BORDER
    line.line.fill.background()


def add_footer(slide: Slide, payload: dict[str, Any]) -> None:
    metadata = payload["metadata"]
    text = (
        f"AMS Applications & Volumetrics Analytics | "
        f"{metadata['customer_name']} | {metadata['project_name']}"
    )
    add_textbox(slide, text, 0.45, 7.14, 11.3, 0.22, size=8, color=MUTED)


def add_panel(
    slide: Slide, x: float, y: float, w: float, h: float, *, fill: RGBColor = WHITE
) -> Any:
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = BORDER
    return shape


def add_section_label(slide: Slide, label: str, x: float, y: float, w: float = 3.0) -> None:
    add_textbox(slide, label.upper(), x, y, w, 0.22, size=8, bold=True, color=MUTED)


def commentary_lookup(payload: dict[str, Any]) -> dict[tuple[str, ...], str]:
    lookup: dict[tuple[str, ...], str] = {}
    for row in payload.get("commentaries", []):
        text = safe_text(row.get("commentary_text")).strip() or strip_html(
            row.get("commentary_html")
        )
        if not text:
            continue
        key = (
            safe_text(row.get("dashboard_area")),
            safe_text(row.get("tab_name")),
            safe_text(row.get("sub_tab_name")),
            safe_text(row.get("section_key")),
            safe_text(row.get("chart_key")),
            safe_text(row.get("scope_filter") or "all"),
            safe_text(row.get("ticket_type_filter") or "all"),
            safe_text(row.get("functional_track_ams_owner") or "all"),
        )
        lookup[key] = text
    return lookup


def normalize_functional_filter(value: str | list[str] | None) -> FunctionalFilter:
    if isinstance(value, list):
        values = sorted({safe_text(item).strip() for item in value if safe_text(item).strip()})
        values = [item for item in values if item.casefold() != "all"]
        if not values:
            return "all"
        if len(values) == 1:
            return values[0]
        return frozenset(values)
    cleaned = safe_text(value).strip()
    return cleaned if cleaned and cleaned.casefold() != "all" else "all"


def functional_filter_values(functional: FunctionalFilter) -> frozenset[str] | None:
    if functional == "all":
        return None
    if isinstance(functional, frozenset):
        return functional
    return frozenset({functional})


def functional_filter_label(functional: FunctionalFilter) -> str:
    values = functional_filter_values(functional)
    if values is None:
        return "all"
    return "; ".join(sorted(values))


def commentary_functional_value(functional: FunctionalFilter) -> str:
    if isinstance(functional, frozenset):
        return "all"
    return functional


def get_commentary(
    lookup: dict[tuple[str, ...], str],
    *,
    dashboard_area: str,
    tab_name: str,
    section_key: str,
    chart_key: str = "",
    sub_tab_name: str = "",
    scope: str = "all",
    ticket_type: str = "all",
    functional: FunctionalFilter = "all",
) -> str:
    functional_key = commentary_functional_value(functional)
    candidates = [
        (
            dashboard_area,
            tab_name,
            sub_tab_name,
            section_key,
            chart_key,
            scope,
            ticket_type,
            functional_key,
        ),
        (
            dashboard_area,
            tab_name,
            sub_tab_name,
            section_key,
            chart_key,
            scope,
            "all",
            functional_key,
        ),
        (
            dashboard_area,
            tab_name,
            sub_tab_name,
            section_key,
            chart_key,
            "all",
            ticket_type,
            functional_key,
        ),
        (
            dashboard_area,
            tab_name,
            sub_tab_name,
            section_key,
            chart_key,
            "all",
            "all",
            functional_key,
        ),
        (dashboard_area, tab_name, sub_tab_name, section_key, chart_key, "all", "all", "all"),
    ]
    for key in candidates:
        if key in lookup:
            return lookup[key]
    return ""


def add_commentary(slide: Slide, text: str, x: float, y: float, w: float, h: float) -> None:
    if not text:
        return
    add_panel(slide, x, y, w, h, fill=RGBColor(245, 249, 252))
    add_textbox(
        slide,
        "Commentary / Inferences",
        x + 0.12,
        y + 0.08,
        w - 0.24,
        0.24,
        size=8,
        bold=True,
        color=MUTED,
    )
    trimmed = text if len(text) <= 520 else f"{text[:517].rstrip()}..."
    add_textbox(slide, trimmed, x + 0.12, y + 0.34, w - 0.24, h - 0.42, size=10, color=TEXT)


def add_table(
    slide: Slide,
    columns: list[str],
    rows: list[list[Any]],
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    font_size: int = 8,
) -> None:
    if not rows:
        add_textbox(slide, "No table data available.", x, y, w, 0.4, size=11, color=MUTED)
        return
    table_shape = slide.shapes.add_table(
        len(rows) + 1, len(columns), Inches(x), Inches(y), Inches(w), Inches(h)
    )
    table = table_shape.table
    for col_index, column in enumerate(columns):
        cell = table.cell(0, col_index)
        cell.text = column
        cell.fill.solid()
        cell.fill.fore_color.rgb = TEAL
        for paragraph in cell.text_frame.paragraphs:
            for run in paragraph.runs:
                set_font(run, size=font_size, bold=True, color=WHITE)
    for row_index, row in enumerate(rows, start=1):
        for col_index, value in enumerate(row):
            cell = table.cell(row_index, col_index)
            cell.text = safe_text(value)
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    set_font(run, size=font_size, color=TEXT)


def add_category_chart(
    slide: Slide,
    title: str,
    categories: list[str],
    series: list[tuple[str, list[float], RGBColor]],
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    chart_type: XL_CHART_TYPE = XL_CHART_TYPE.COLUMN_CLUSTERED,
    legend: bool = True,
) -> None:
    add_textbox(slide, title, x, y, w, 0.28, size=12, bold=True, color=TEXT)
    if not categories or not series:
        add_textbox(slide, "No chart data available.", x, y + 0.35, w, 0.35, size=10, color=MUTED)
        return
    chart_data = CategoryChartData()
    chart_data.categories = categories
    for name, values, _color in series:
        chart_data.add_series(name, values)
    chart_shape = slide.shapes.add_chart(
        chart_type,
        Inches(x),
        Inches(y + 0.35),
        Inches(w),
        Inches(h - 0.35),
        chart_data,
    )
    chart = chart_shape.chart
    chart.has_legend = legend
    if legend:
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
    chart.has_title = False
    if chart.plots:
        plot = chart.plots[0]
        plot.has_data_labels = True
        plot.data_labels.number_format = "#,##0.0"
        plot.data_labels.position = XL_LABEL_POSITION.OUTSIDE_END
    for series_index, chart_series in enumerate(chart.series):
        chart_series.format.fill.solid()
        chart_series.format.fill.fore_color.rgb = series[series_index][2]


def add_pie_chart(
    slide: Slide,
    title: str,
    points: list[tuple[str, float]],
    x: float,
    y: float,
    w: float,
    h: float,
) -> None:
    add_textbox(slide, title, x, y, w, 0.34, size=11, bold=True, color=TEXT)
    filtered = [(label, value) for label, value in points if value > 0]
    if not filtered:
        add_textbox(slide, "No chart data available.", x, y + 0.42, w, 0.35, size=10, color=MUTED)
        return
    chart_data = CategoryChartData()
    chart_data.categories = [label for label, _value in filtered]
    chart_data.add_series("Average monthly volume", [value for _label, value in filtered])
    chart_shape = slide.shapes.add_chart(
        XL_CHART_TYPE.PIE,
        Inches(x),
        Inches(y + 0.35),
        Inches(w),
        Inches(h - 0.35),
        chart_data,
    )
    chart = chart_shape.chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False
    plot = chart.plots[0]
    plot.has_data_labels = True
    plot.data_labels.show_percentage = True
    plot.data_labels.number_format = "0.0%"


def add_horizontal_bars(
    slide: Slide,
    title: str,
    points: list[tuple[str, float, str]],
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    color: RGBColor = TEAL,
) -> None:
    add_textbox(slide, title, x, y, w, 0.34, size=12, bold=True, color=TEXT)
    if not points:
        add_textbox(slide, "No chart data available.", x, y + 0.42, w, 0.35, size=10, color=MUTED)
        return
    max_value = max(max(value for _label, value, _display in points), 1)
    row_h = (h - 0.5) / max(len(points), 1)
    bar_x = x + min(2.8, w * 0.38)
    bar_w = w - (bar_x - x) - 0.7
    for index, (label, value, display) in enumerate(points):
        row_y = y + 0.48 + index * row_h
        add_textbox(
            slide, label[:42], x, row_y, bar_x - x - 0.08, min(row_h, 0.3), size=7, color=TEXT
        )
        width = bar_w * (value / max_value)
        rect = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(bar_x),
            Inches(row_y + 0.03),
            Inches(max(width, 0.03)),
            Inches(max(row_h * 0.42, 0.08)),
        )
        rect.fill.solid()
        rect.fill.fore_color.rgb = color
        rect.line.fill.background()
        add_textbox(
            slide,
            display,
            bar_x + width + 0.04,
            row_y,
            0.6,
            min(row_h, 0.3),
            size=7,
            bold=True,
            color=TEXT,
        )


def normalize_ticket_type(value: str) -> str:
    normalized = value.strip().casefold()
    if normalized in {"incident", "incidents", "inc"}:
        return "incident"
    if normalized in {"sc_task", "sc tasks", "service_catalog_task", "service catalog task"}:
        return "sc_task"
    return "all"


def normalize_scope(value: str) -> str:
    normalized = value.strip().casefold()
    if normalized in {"in_scope", "in-scope", "in scope"}:
        return "in_scope"
    if normalized in {"out_of_scope", "out-of-scope", "out of scope"}:
        return "out_of_scope"
    return "all"


def row_matches(
    row: dict[str, Any],
    *,
    scope: str,
    ticket_type: str,
    functional: FunctionalFilter,
) -> bool:
    functional_values = functional_filter_values(functional)
    return (
        (scope == "all" or row.get("scope") == scope)
        and (ticket_type == "all" or row.get("ticket_type") == ticket_type)
        and (
            functional_values is None
            or safe_text(row.get("functional_track_ams_owner")) in functional_values
        )
    )


def app_matches(row: dict[str, Any], *, functional: FunctionalFilter) -> bool:
    functional_values = functional_filter_values(functional)
    if functional_values is None:
        return True
    return safe_text(row.get("functional_track_ams_owner")) in functional_values


def aggregate_periods(
    payload: dict[str, Any], *, scope: str, ticket_type: str, functional: str
) -> list[dict[str, Any]]:
    rows = [
        row
        for row in payload["volumetrics"]["monthly_rows"]
        if row_matches(row, scope=scope, ticket_type=ticket_type, functional=functional)
    ]
    result = []
    for period in payload["volumetrics"]["periods"]:
        matching = [row for row in rows if row["period_key"] == period["period_key"]]
        result.append(
            {
                "label": period["period_label"],
                "created": sum(int(row.get("created_count") or 0) for row in matching),
                "resolved": sum(int(row.get("resolved_closed_count") or 0) for row in matching),
                "canceled": sum(
                    int(row.get("canceled_closed_incomplete_count") or 0) for row in matching
                ),
                "backlog": sum(int(row.get("backlog_open") or 0) for row in matching),
            },
        )
    return result


def top_active_users_points(
    app_rows: list[dict[str, Any]], *, functional: str, limit: int = 10
) -> list[tuple[str, float, str]]:
    by_parent: dict[str, int] = {}
    for row in app_rows:
        if not app_matches(row, functional=functional):
            continue
        users = int(row.get("active_users") or 0)
        parent = safe_text(row.get("parent_application_name")).strip()
        if users <= 0 or not parent:
            continue
        by_parent[parent] = max(by_parent.get(parent, 0), users)
    return [
        (label, float(users), fmt_number(users))
        for label, users in sorted(by_parent.items(), key=lambda item: (-item[1], item[0]))[:limit]
    ]


def count_by(
    app_rows: list[dict[str, Any]], field: str, *, functional: str
) -> list[tuple[str, float]]:
    counts: dict[str, int] = {}
    for row in app_rows:
        if not app_matches(row, functional=functional):
            continue
        label = safe_text(row.get(field)).strip() or "(blank)"
        counts[label] = counts.get(label, 0) + 1
    return [
        (label, float(count))
        for label, count in sorted(counts.items(), key=lambda item: (-item[1], item[0]))
    ]


def top_volume_points(
    payload: dict[str, Any],
    *,
    scope: str,
    ticket_type: str,
    functional: str,
    limit: int = 10,
) -> list[tuple[str, float, str]]:
    rows = [
        row
        for row in payload["volumetrics"]["detailed_volume_trends"]["application_rows"]
        if row_matches(row, scope=scope, ticket_type=ticket_type, functional=functional)
    ]
    by_app: dict[str, int] = {}
    for row in rows:
        app_name = safe_text(row.get("application_name")) or "(blank)"
        by_app[app_name] = by_app.get(app_name, 0) + int(row.get("created_count") or 0)
    points = []
    for label, total in by_app.items():
        average = total / 6
        if average > 0:
            points.append((label, average, fmt_number(average, 1)))
    return sorted(points, key=lambda item: (-item[1], item[0]))[:limit]


def tickets_per_user_points(
    payload: dict[str, Any],
    *,
    scope: str,
    ticket_type: str,
    functional: str,
    limit: int = 10,
) -> list[tuple[str, float, str]]:
    active_users = {
        safe_text(row.get("business_service_ci_name")): int(row.get("active_users") or 0)
        for row in payload["applications"]["rows"]
        if app_matches(row, functional=functional) and int(row.get("active_users") or 0) > 0
    }
    totals: dict[str, int] = {}
    for row in payload["volumetrics"]["detailed_volume_trends"]["application_rows"]:
        if not row_matches(row, scope=scope, ticket_type=ticket_type, functional=functional):
            continue
        app_name = safe_text(row.get("application_name"))
        if app_name not in active_users:
            continue
        totals[app_name] = totals.get(app_name, 0) + int(row.get("created_count") or 0)
    points = []
    for label, total in totals.items():
        users = active_users.get(label, 0)
        if users <= 0:
            continue
        ratio = (total / 6) / users
        points.append((label, ratio, fmt_number(ratio, 2)))
    return sorted(points, key=lambda item: (-item[1], item[0]))[:limit]


def distribution_points(
    payload: dict[str, Any],
    field: str,
    selected_type: str,
    *,
    scope: str,
    functional: str,
) -> list[tuple[str, float]]:
    totals: dict[str, int] = {}
    for row in payload["volumetrics"]["detailed_volume_trends"]["application_rows"]:
        if not row_matches(row, scope=scope, ticket_type=selected_type, functional=functional):
            continue
        label = safe_text(row.get(field)).strip() or "(blank)"
        totals[label] = totals.get(label, 0) + int(row.get("created_count") or 0)
    return [
        (label, count / 6)
        for label, count in sorted(totals.items(), key=lambda item: (-item[1], item[0]))
        if count > 0
    ]


def priority_periods(
    payload: dict[str, Any], *, scope: str, ticket_type: str, functional: str
) -> tuple[list[str], list[dict[str, Any]]]:
    source = payload["volumetrics"]["overall_volume_trends"]["priority_distribution"]
    priorities = list(source.get("priorities") or [])
    rows = [
        row
        for row in source.get("rows", [])
        if row_matches(row, scope=scope, ticket_type=ticket_type, functional=functional)
    ]
    period_rows = []
    for period in payload["volumetrics"]["periods"]:
        values = {priority: 0 for priority in priorities}
        for row in rows:
            if row.get("period_key") == period["period_key"]:
                values[row["priority"]] = values.get(row["priority"], 0) + int(
                    row.get("ticket_count") or 0
                )
        period_rows.append({"label": period["period_label"], **values})
    return priorities, period_rows


def sla_points(
    payload: dict[str, Any], kind: str, *, scope: str, functional: str
) -> list[tuple[str, float]]:
    rows = [
        row
        for row in payload["volumetrics"]["overall_sla_trends"].get("rows", [])
        if row_matches(row, scope=scope, ticket_type="incident", functional=functional)
    ]
    points: list[tuple[str, float]] = []
    for period in payload["volumetrics"]["periods"]:
        matching = [row for row in rows if row["period_key"] == period["period_key"]]
        captured = sum(int(row.get(f"{kind}_sla_captured_count") or 0) for row in matching)
        adhered = sum(int(row.get(f"{kind}_sla_adhered_count") or 0) for row in matching)
        points.append((period["period_label"], (adhered / captured * 100) if captured else 0))
    return points


def mttr_points(
    payload: dict[str, Any],
    selected_type: str,
    priority: str,
    *,
    scope: str,
    functional: str,
) -> list[float]:
    rows = [
        row
        for row in payload["volumetrics"]["kpi_trends"]["mttr"].get("rows", [])
        if row_matches(row, scope=scope, ticket_type=selected_type, functional=functional)
        and row.get("priority") == priority
    ]
    values_by_period: dict[str, tuple[float, int]] = {}
    for row in rows:
        current_sum, current_count = values_by_period.get(row["period_key"], (0.0, 0))
        values_by_period[row["period_key"]] = (
            current_sum + float(row.get("business_duration_seconds_sum") or 0),
            current_count + int(row.get("ticket_count") or 0),
        )
    values = []
    for period in payload["volumetrics"]["periods"]:
        seconds, count = values_by_period.get(period["period_key"], (0.0, 0))
        values.append(seconds / count / 86400 if count else 0)
    return values


def duration_bucket_rows(
    payload: dict[str, Any],
    selected_type: str,
    *,
    scope: str,
    functional: str,
) -> list[dict[str, Any]]:
    source = payload["volumetrics"]["kpi_trends"]["duration_buckets"]
    rows = [
        row
        for row in source.get("rows", [])
        if row_matches(row, scope=scope, ticket_type=selected_type, functional=functional)
    ]
    result = []
    for period in source.get("periods", []):
        values = {bucket: 0 for bucket in source.get("buckets", DURATION_BUCKETS)}
        for row in rows:
            if row.get("period_key") == period["period_key"]:
                values[row["bucket"]] += int(row.get("ticket_count") or 0)
        result.append({"label": period["period_label"], **values})
    return result


def add_title_slide(
    presentation: PresentationType,
    payload: dict[str, Any],
    *,
    scope: str,
    ticket_type: str,
    functional: FunctionalFilter,
) -> None:
    slide = add_blank_slide(presentation)
    add_panel(slide, 0, 0, 13.33, 7.5, fill=RGBColor(236, 253, 245))
    add_textbox(
        slide,
        "AMS Applications & Volumetrics Analytics",
        0.8,
        1.65,
        11.7,
        0.8,
        size=34,
        bold=True,
        color=TEAL,
        align=PP_ALIGN.CENTER,
    )
    metadata = payload["metadata"]
    add_textbox(
        slide,
        f"{metadata['customer_name']} | {metadata['project_name']}",
        1.0,
        2.55,
        11.3,
        0.35,
        size=16,
        bold=True,
        color=TEXT,
        align=PP_ALIGN.CENTER,
    )
    complete_range = (
        f"Complete-month range: {metadata.get('complete_month_from') or 'N/A'} to "
        f"{metadata.get('complete_month_to') or 'N/A'}"
    )
    add_textbox(
        slide, complete_range, 1.0, 3.05, 11.3, 0.3, size=12, color=MUTED, align=PP_ALIGN.CENTER
    )
    add_textbox(
        slide,
        "Scope: "
        f"{scope} | Ticket Type: {ticket_type} | Functional Track / AMS Owner: "
        f"{functional_filter_label(functional)}",
        1.0,
        3.45,
        11.3,
        0.3,
        size=10,
        color=MUTED,
        align=PP_ALIGN.CENTER,
    )


def add_overview_slide(
    presentation: PresentationType, payload: dict[str, Any], commentary: str
) -> None:
    slide = add_blank_slide(presentation)
    add_header(slide, "Overview", "Summary tiles use complete-month dashboard data.")
    overview = payload["overview"]
    inventory = overview["application_inventory"]
    tickets = overview["tickets"]
    tiles = [
        ("Applications", fmt_number(inventory.get("total_applications")), "Application Inventory"),
        (
            "Functional Tracks",
            fmt_number(inventory.get("functional_track_count")),
            "Inventory grouping",
        ),
        (
            "In-Scope Tickets",
            fmt_number(tickets.get("total_in_scope_tickets")),
            "Tickets table only",
        ),
        ("Incidents", fmt_number(tickets.get("incident_count")), "In-scope Incident volume"),
        ("SC Tasks", fmt_number(tickets.get("sc_task_count")), "In-scope SC Task volume"),
        (
            "Apps Driving 80%",
            fmt_number(tickets.get("applications_80pct_monthly_volume_count")),
            "Average monthly volume",
        ),
    ]
    for index, (label, value, helper) in enumerate(tiles):
        col = index % 3
        row = index // 3
        x = 0.65 + col * 4.05
        y = 1.35 + row * 1.25
        add_panel(slide, x, y, 3.75, 0.95, fill=WHITE)
        add_textbox(slide, label, x + 0.15, y + 0.1, 3.4, 0.22, size=8, bold=True, color=MUTED)
        add_textbox(slide, value, x + 0.15, y + 0.34, 3.4, 0.33, size=20, bold=True, color=TEAL)
        add_textbox(slide, helper, x + 0.15, y + 0.7, 3.4, 0.18, size=8, color=MUTED)
    add_commentary(slide, commentary, 0.65, 4.05, 12.0, 1.25)
    add_footer(slide, payload)


def add_applications_summary_slide(
    presentation: PresentationType,
    payload: dict[str, Any],
    lookup: dict[tuple[str, ...], str],
    functional: str,
) -> None:
    slide = add_blank_slide(presentation)
    add_header(slide, "Applications Summary", "Application Inventory charts.")
    app_rows = payload["applications"]["rows"]
    chart_specs = [
        ("Strategic", "strategic", 0.55, 1.15),
        ("Lifecycle Stage", "lifecycle_stage_status", 4.55, 1.15),
        ("Hosting Env", "hosting_env", 8.55, 1.15),
        ("Architecture Type", "architecture_type", 0.55, 3.85),
        ("Install Type", "install_type", 4.55, 3.85),
    ]
    for title, field, x, y in chart_specs:
        add_pie_chart(
            slide, title, count_by(app_rows, field, functional=functional)[:8], x, y, 3.65, 2.15
        )
        text = get_commentary(
            lookup,
            dashboard_area="applications",
            tab_name="applications",
            section_key="applications_charts",
            chart_key={"lifecycle_stage_status": "lifecycle_stage"}.get(field, field),
            functional=functional,
        )
        if text:
            add_commentary(slide, text, x, y + 2.15, 3.65, 0.45)
    summary_text = get_commentary(
        lookup,
        dashboard_area="applications",
        tab_name="applications",
        section_key="applications_summary",
        functional=functional,
    )
    add_commentary(slide, summary_text, 8.55, 4.0, 3.9, 1.9)
    add_footer(slide, payload)


def add_active_users_slide(
    presentation: PresentationType, payload: dict[str, Any], commentary: str, *, functional: str
) -> None:
    slide = add_blank_slide(presentation)
    add_header(
        slide,
        "Top Parent Business Applications by Active Users",
        "One row per parent application; duplicates use highest Active Users.",
    )
    add_horizontal_bars(
        slide,
        "Top 10 Active Users",
        top_active_users_points(payload["applications"]["rows"], functional=functional),
        0.55,
        1.25,
        8.0,
        4.8,
        color=TEAL,
    )
    add_commentary(slide, commentary, 8.8, 1.45, 3.9, 2.2)
    add_footer(slide, payload)


def add_application_table_slides(
    presentation: PresentationType, payload: dict[str, Any], commentary: str, *, functional: str
) -> None:
    rows = [
        row for row in payload["applications"]["rows"] if app_matches(row, functional=functional)
    ]
    columns = [
        ("Application", "business_service_ci_name"),
        ("Parent App", "parent_application_name"),
        ("SAP", "sap_non_sap"),
        ("Functional / Owner", "functional_track_ams_owner"),
        ("Active Users", "active_users"),
    ]
    chunk_size = 12
    chunks = [rows[index : index + chunk_size] for index in range(0, len(rows), chunk_size)] or [[]]
    for index, chunk in enumerate(chunks, start=1):
        slide = add_blank_slide(presentation)
        suffix = f" ({index}/{len(chunks)})" if len(chunks) > 1 else ""
        add_header(
            slide, f"Application List{suffix}", f"{len(rows)} inventory rows in selected context."
        )
        add_table(
            slide,
            [label for label, _field in columns],
            [
                [
                    fmt_number(row.get(field)) if field == "active_users" else row.get(field, "")
                    for _label, field in columns
                ]
                for row in chunk
            ],
            0.45,
            1.25,
            12.35,
            4.5,
            font_size=7,
        )
        if index == 1:
            add_commentary(slide, commentary, 0.45, 5.9, 12.35, 0.9)
        add_footer(slide, payload)


def add_overall_volume_slides(
    presentation: PresentationType,
    payload: dict[str, Any],
    lookup: dict[tuple[str, ...], str],
    *,
    scope: str,
    ticket_type: str,
    functional: str,
) -> None:
    periods = aggregate_periods(
        payload, scope=scope, ticket_type=ticket_type, functional=functional
    )
    labels = [row["label"] for row in periods]
    slide = add_blank_slide(presentation)
    add_header(
        slide,
        "Overall Volume Trends - Created / Completed / Backlog",
        "Monthly complete-period trend.",
    )
    add_category_chart(
        slide,
        "Created vs Resolved/Closed vs Canceled / Closed Incomplete",
        labels,
        [
            ("Created", [row["created"] for row in periods], TEAL),
            ("Resolved/Closed", [row["resolved"] for row in periods], BLUE),
            ("Canceled / Closed Incomplete", [row["canceled"] for row in periods], RED),
        ],
        0.55,
        1.25,
        7.2,
        3.4,
    )
    add_category_chart(
        slide,
        "Backlog(Open)",
        labels,
        [("Backlog(Open)", [row["backlog"] for row in periods], ORANGE)],
        7.95,
        1.25,
        4.8,
        3.4,
        chart_type=XL_CHART_TYPE.LINE_MARKERS,
    )
    add_commentary(
        slide,
        get_commentary(
            lookup,
            dashboard_area="volumetrics",
            tab_name="volumetrics_sla",
            sub_tab_name="overall_volume_trends",
            section_key="overall_volume_trends",
            chart_key="created_resolved_canceled",
            scope=scope,
            ticket_type=ticket_type,
            functional=functional,
        ),
        0.55,
        5.0,
        6.0,
        0.9,
    )
    add_commentary(
        slide,
        get_commentary(
            lookup,
            dashboard_area="volumetrics",
            tab_name="volumetrics_sla",
            sub_tab_name="overall_volume_trends",
            section_key="overall_volume_trends",
            chart_key="backlog",
            scope=scope,
            ticket_type=ticket_type,
            functional=functional,
        ),
        6.8,
        5.0,
        5.95,
        0.9,
    )
    add_footer(slide, payload)

    slide = add_blank_slide(presentation)
    add_header(
        slide,
        "Overall Volume Trends - Patterns and Priority",
        "Created patterns, hourly split, and priority distribution.",
    )
    priorities, priority_rows = priority_periods(
        payload, scope=scope, ticket_type=ticket_type, functional=functional
    )
    add_table(
        slide,
        ["Period", *priorities],
        [
            [row.get("label"), *[fmt_number(row.get(priority)) for priority in priorities]]
            for row in priority_rows
        ],
        0.55,
        1.25,
        5.8,
        2.2,
        font_size=7,
    )
    add_category_chart(
        slide,
        "Priority-wise ticket distribution",
        [row["label"] for row in priority_rows],
        [
            (
                priority,
                [row.get(priority, 0) for row in priority_rows],
                CHART_COLORS[index % len(CHART_COLORS)],
            )
            for index, priority in enumerate(priorities)
        ],
        6.65,
        1.25,
        6.0,
        3.2,
        chart_type=XL_CHART_TYPE.COLUMN_STACKED,
    )
    add_commentary(
        slide,
        get_commentary(
            lookup,
            dashboard_area="volumetrics",
            tab_name="volumetrics_sla",
            sub_tab_name="overall_volume_trends",
            section_key="overall_volume_trends",
            chart_key="priority_distribution",
            scope=scope,
            ticket_type=ticket_type,
            functional=functional,
        ),
        0.55,
        4.85,
        12.1,
        0.95,
    )
    add_footer(slide, payload)


def add_sla_slide(
    presentation: PresentationType,
    payload: dict[str, Any],
    lookup: dict[tuple[str, ...], str],
    kind: str,
    title: str,
    chart_key: str,
    *,
    scope: str,
    commentary_ticket_type: str,
    functional: str,
) -> None:
    slide = add_blank_slide(presentation)
    add_header(slide, title, "Incident SLA trend; SC Tasks are excluded from SLA calculations.")
    points = sla_points(payload, kind, scope=scope, functional=functional)
    add_category_chart(
        slide,
        f"{title} %",
        [label for label, _value in points],
        [(title, [value for _label, value in points], TEAL if kind == "response" else BLUE)],
        0.65,
        1.25,
        8.1,
        3.6,
        chart_type=XL_CHART_TYPE.LINE_MARKERS,
    )
    add_table(
        slide,
        ["Period", "Adherence %"],
        [[label, fmt_pct(value)] for label, value in points],
        9.1,
        1.35,
        3.3,
        3.2,
        font_size=8,
    )
    add_commentary(
        slide,
        get_commentary(
            lookup,
            dashboard_area="volumetrics",
            tab_name="volumetrics_sla",
            sub_tab_name="overall_sla_trends",
            section_key="overall_sla_trends",
            chart_key=chart_key,
            scope=scope,
            ticket_type=commentary_ticket_type,
            functional=functional,
        ),
        0.65,
        5.2,
        11.8,
        0.85,
    )
    add_footer(slide, payload)


def add_detailed_volume_slides(
    presentation: PresentationType,
    payload: dict[str, Any],
    lookup: dict[tuple[str, ...], str],
    *,
    scope: str,
    ticket_type: str,
    functional: str,
) -> None:
    slide = add_blank_slide(presentation)
    add_header(
        slide,
        "Detailed Volume Trends - Top Applications",
        "Latest complete six-month average monthly volume.",
    )
    add_horizontal_bars(
        slide,
        "Top High-Volume Applications",
        top_volume_points(payload, scope=scope, ticket_type=ticket_type, functional=functional),
        0.55,
        1.25,
        5.8,
        4.1,
        color=TEAL,
    )
    add_horizontal_bars(
        slide,
        "Tickets per User per Month by Application",
        tickets_per_user_points(
            payload, scope=scope, ticket_type=ticket_type, functional=functional
        ),
        6.85,
        1.25,
        5.8,
        4.1,
        color=PURPLE,
    )
    add_commentary(
        slide,
        get_commentary(
            lookup,
            dashboard_area="volumetrics",
            tab_name="volumetrics_sla",
            sub_tab_name="detailed_volume_trends",
            section_key="detailed_volume_trends",
            chart_key="top_high_volume_applications",
            scope=scope,
            ticket_type=ticket_type,
            functional=functional,
        ),
        0.55,
        5.65,
        5.8,
        0.75,
    )
    add_commentary(
        slide,
        get_commentary(
            lookup,
            dashboard_area="volumetrics",
            tab_name="volumetrics_sla",
            sub_tab_name="detailed_volume_trends",
            section_key="detailed_volume_trends",
            chart_key="tickets_per_user_application",
            scope=scope,
            ticket_type=ticket_type,
            functional=functional,
        ),
        6.85,
        5.65,
        5.8,
        0.75,
    )
    add_footer(slide, payload)

    slide = add_blank_slide(presentation)
    add_header(
        slide,
        "Detailed Volume Trends - Distribution Splits",
        "Average monthly created ticket volume.",
    )
    rows = [
        ("SAP / Non-SAP", "sap_non_sap", "sap_non_sap_distribution_row", 1.1),
        ("Architecture Type", "architecture_type", "architecture_type_distribution_row", 3.05),
        ("Install Type", "install_type", "install_type_distribution_row", 5.0),
    ]
    for row_title, field, key, y in rows:
        add_section_label(slide, row_title, 0.55, y - 0.15)
        add_pie_chart(
            slide,
            "Tickets",
            distribution_points(payload, field, "all", scope=scope, functional=functional),
            0.55,
            y,
            3.6,
            1.75,
        )
        add_pie_chart(
            slide,
            "Incidents",
            distribution_points(payload, field, "incident", scope=scope, functional=functional),
            4.55,
            y,
            3.6,
            1.75,
        )
        add_pie_chart(
            slide,
            "SC Tasks",
            distribution_points(payload, field, "sc_task", scope=scope, functional=functional),
            8.55,
            y,
            3.6,
            1.75,
        )
        note = get_commentary(
            lookup,
            dashboard_area="volumetrics",
            tab_name="volumetrics_sla",
            sub_tab_name="detailed_volume_trends",
            section_key="detailed_volume_trends",
            chart_key=key,
            scope=scope,
            ticket_type=ticket_type,
            functional=functional,
        )
        if note and y == 5.0:
            add_commentary(slide, note, 0.55, 6.55, 12.1, 0.45)
    add_footer(slide, payload)

    slide = add_blank_slide(presentation)
    add_header(
        slide,
        "Detailed Volume Trends - Hosting Env",
        "Average monthly created ticket volume by Hosting Env.",
    )
    add_pie_chart(
        slide,
        "Tickets",
        distribution_points(payload, "hosting_env", "all", scope=scope, functional=functional),
        0.55,
        1.55,
        3.6,
        2.8,
    )
    add_pie_chart(
        slide,
        "Incidents",
        distribution_points(payload, "hosting_env", "incident", scope=scope, functional=functional),
        4.55,
        1.55,
        3.6,
        2.8,
    )
    add_pie_chart(
        slide,
        "SC Tasks",
        distribution_points(payload, "hosting_env", "sc_task", scope=scope, functional=functional),
        8.55,
        1.55,
        3.6,
        2.8,
    )
    add_commentary(
        slide,
        get_commentary(
            lookup,
            dashboard_area="volumetrics",
            tab_name="volumetrics_sla",
            sub_tab_name="detailed_volume_trends",
            section_key="detailed_volume_trends",
            chart_key="hosting_env_distribution_row",
            scope=scope,
            ticket_type=ticket_type,
            functional=functional,
        ),
        0.55,
        5.25,
        12.1,
        0.9,
    )
    add_footer(slide, payload)


def add_mttr_slide(
    presentation: PresentationType,
    payload: dict[str, Any],
    lookup: dict[tuple[str, ...], str],
    ticket_type: str,
    title: str,
    *,
    scope: str,
    commentary_ticket_type: str,
    functional: str,
) -> None:
    slide = add_blank_slide(presentation)
    add_header(slide, title, "MTTR uses average business duration seconds converted to days.")
    labels = [period["period_label"] for period in payload["volumetrics"]["periods"]]
    prefix = "Incident" if ticket_type == "incident" else "SC Task"
    for index, pair in enumerate((("P1", "P2"), ("P3", "P4"))):
        x = 0.65
        y = 1.25 + index * 2.55
        add_category_chart(
            slide,
            f"{prefix} {' / '.join(pair)} MTTR",
            labels,
            [
                (
                    priority,
                    mttr_points(payload, ticket_type, priority, scope=scope, functional=functional),
                    PRIORITY_COLORS[MTTR_PRIORITIES.index(priority)],
                )
                for priority in pair
            ],
            x,
            y,
            7.3,
            2.25,
            chart_type=XL_CHART_TYPE.LINE_MARKERS,
        )
        key = f"{ticket_type}_{'_'.join(priority.lower() for priority in pair)}_mttr"
        add_commentary(
            slide,
            get_commentary(
                lookup,
                dashboard_area="volumetrics",
                tab_name="volumetrics_sla",
                sub_tab_name="kpi_trends",
                section_key="kpi_trends",
                chart_key=key,
                scope=scope,
                ticket_type=commentary_ticket_type,
                functional=functional,
            ),
            8.25,
            y + 0.35,
            4.1,
            1.15,
        )
    add_footer(slide, payload)


def add_duration_slide(
    presentation: PresentationType,
    payload: dict[str, Any],
    lookup: dict[tuple[str, ...], str],
    *,
    scope: str,
    commentary_ticket_type: str,
    functional: str,
) -> None:
    slide = add_blank_slide(presentation)
    add_header(slide, "KPI Trends - Duration Buckets", "Latest three complete months.")
    specs = [
        ("Incident Resolved Duration", "incident", "incident_duration_buckets_row", 1.25),
        ("SC Task Closed Duration", "sc_task", "sc_task_duration_buckets_row", 4.05),
    ]
    for title, selected_type, key, y in specs:
        rows = duration_bucket_rows(payload, selected_type, scope=scope, functional=functional)
        add_category_chart(
            slide,
            title,
            [row["label"] for row in rows],
            [
                (
                    bucket,
                    [row.get(bucket, 0) for row in rows],
                    CHART_COLORS[index % len(CHART_COLORS)],
                )
                for index, bucket in enumerate(DURATION_BUCKETS)
            ],
            0.65,
            y,
            7.5,
            2.2,
            chart_type=XL_CHART_TYPE.COLUMN_CLUSTERED,
        )
        add_commentary(
            slide,
            get_commentary(
                lookup,
                dashboard_area="volumetrics",
                tab_name="volumetrics_sla",
                sub_tab_name="kpi_trends",
                section_key="kpi_trends",
                chart_key=key,
                scope=scope,
                ticket_type=commentary_ticket_type,
                functional=functional,
            ),
            8.45,
            y + 0.4,
            3.95,
            1.1,
        )
    add_footer(slide, payload)


def add_placeholder_slide(
    presentation: PresentationType, payload: dict[str, Any], title: str, body: str
) -> None:
    slide = add_blank_slide(presentation)
    add_header(slide, title)
    add_panel(slide, 1.1, 2.2, 11.0, 2.0, fill=LIGHT_BG)
    add_textbox(
        slide, body, 1.35, 2.65, 10.5, 0.7, size=18, bold=True, color=MUTED, align=PP_ALIGN.CENTER
    )
    add_footer(slide, payload)


def build_powerpoint_export(
    db: Session,
    project_id: UUID,
    *,
    scope_filter: str = "in_scope",
    ticket_type_filter: str = "all",
    functional_track_ams_owner: str | list[str] = "all",
    include_commentary: bool = True,
) -> tuple[bytes, str]:
    payload = build_offline_dashboard_payload(db, project_id)
    metadata = payload["metadata"]
    scope = normalize_scope(scope_filter)
    ticket_type = normalize_ticket_type(ticket_type_filter)
    functional = normalize_functional_filter(functional_track_ams_owner)
    lookup = commentary_lookup(payload) if include_commentary else {}

    presentation = Presentation()
    presentation.slide_width = Inches(13.333)
    presentation.slide_height = Inches(7.5)

    add_title_slide(
        presentation, payload, scope=scope, ticket_type=ticket_type, functional=functional
    )
    add_overview_slide(
        presentation,
        payload,
        get_commentary(
            lookup,
            dashboard_area="dashboard",
            tab_name="overview",
            section_key="overview_summary",
        ),
    )
    add_applications_summary_slide(presentation, payload, lookup, functional=functional)
    add_active_users_slide(
        presentation,
        payload,
        get_commentary(
            lookup,
            dashboard_area="applications",
            tab_name="applications",
            section_key="applications_charts",
            chart_key="top_active_users",
            functional=functional,
        ),
        functional=functional,
    )
    add_application_table_slides(
        presentation,
        payload,
        get_commentary(
            lookup,
            dashboard_area="applications",
            tab_name="applications",
            section_key="application_list",
            chart_key="application_list",
            functional=functional,
        ),
        functional=functional,
    )
    add_overall_volume_slides(
        presentation,
        payload,
        lookup,
        scope=scope,
        ticket_type=ticket_type,
        functional=functional,
    )
    add_sla_slide(
        presentation,
        payload,
        lookup,
        "response",
        "Overall SLA Trends - Response SLA",
        "response_sla_adherence",
        scope=scope,
        commentary_ticket_type=ticket_type,
        functional=functional,
    )
    add_sla_slide(
        presentation,
        payload,
        lookup,
        "resolution",
        "Overall SLA Trends - Resolution SLA",
        "resolution_sla_adherence",
        scope=scope,
        commentary_ticket_type=ticket_type,
        functional=functional,
    )
    add_detailed_volume_slides(
        presentation,
        payload,
        lookup,
        scope=scope,
        ticket_type=ticket_type,
        functional=functional,
    )
    add_mttr_slide(
        presentation,
        payload,
        lookup,
        "incident",
        "KPI Trends - Incident MTTR",
        scope=scope,
        commentary_ticket_type=ticket_type,
        functional=functional,
    )
    add_mttr_slide(
        presentation,
        payload,
        lookup,
        "sc_task",
        "KPI Trends - SC Task MTTR",
        scope=scope,
        commentary_ticket_type=ticket_type,
        functional=functional,
    )
    add_duration_slide(
        presentation,
        payload,
        lookup,
        scope=scope,
        commentary_ticket_type=ticket_type,
        functional=functional,
    )
    add_placeholder_slide(
        presentation,
        payload,
        "Performance Trends",
        "Performance Trends will be added in a future dashboard prompt.",
    )
    add_placeholder_slide(
        presentation,
        payload,
        "Category-wise Trends",
        payload["volumetrics"]
        .get("placeholders", {})
        .get(
            "category_wise_trends",
            "Detailed requirements for this section will be added in a future prompt.",
        ),
    )

    output = BytesIO()
    presentation.save(output)
    exported_at = metadata.get("exported_at")
    filename = powerpoint_filename(
        exported_at if isinstance(exported_at, datetime) else datetime.now(UTC)
    )
    return output.getvalue(), filename
